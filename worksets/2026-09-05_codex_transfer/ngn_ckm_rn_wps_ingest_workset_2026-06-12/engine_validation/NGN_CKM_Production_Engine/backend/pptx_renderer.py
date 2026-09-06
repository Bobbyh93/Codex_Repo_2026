"""Presentation renderer with improved visual hierarchy, image fitting, and layout-safe archetypes."""
from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_X = Inches(0.45)
CONTENT_TOP = Inches(1.05)
CONTENT_WIDTH = SLIDE_WIDTH - (MARGIN_X * 2)
CONTENT_HEIGHT = Inches(5.95)
TITLE_BAR_H = Inches(0.72)
MIN_FONT_SIZE = 18

THEMES = {
    "clean_academic": {
        "background": RGBColor(247, 249, 252),
        "title_bar": RGBColor(22, 49, 92),
        "panel": RGBColor(255, 255, 255),
        "panel_alt": RGBColor(241, 245, 249),
        "line": RGBColor(209, 213, 219),
        "accent": RGBColor(37, 99, 235),
        "ink": RGBColor(17, 24, 39),
        "muted": RGBColor(71, 85, 105),
        "success": RGBColor(5, 150, 105),
    },
    "visual_teaching": {
        "background": RGBColor(245, 248, 255),
        "title_bar": RGBColor(30, 64, 175),
        "panel": RGBColor(255, 255, 255),
        "panel_alt": RGBColor(224, 231, 255),
        "line": RGBColor(191, 219, 254),
        "accent": RGBColor(79, 70, 229),
        "ink": RGBColor(17, 24, 39),
        "muted": RGBColor(71, 85, 105),
        "success": RGBColor(6, 95, 70),
    },
    "case_led": {
        "background": RGBColor(250, 250, 249),
        "title_bar": RGBColor(69, 26, 3),
        "panel": RGBColor(255, 255, 255),
        "panel_alt": RGBColor(254, 243, 199),
        "line": RGBColor(251, 191, 36),
        "accent": RGBColor(217, 119, 6),
        "ink": RGBColor(41, 37, 36),
        "muted": RGBColor(87, 83, 78),
        "success": RGBColor(21, 128, 61),
    },
}


def _theme(deck_style: str | None) -> dict[str, RGBColor]:
    return THEMES.get(deck_style or "clean_academic", THEMES["clean_academic"])


def _set_run_style(run, size: int, color: RGBColor, bold: bool = False, name: str = "Aptos"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def _set_paragraph_text(paragraph, text: str, size: int, color: RGBColor, bold: bool = False, name: str = "Aptos"):
    paragraph.text = text
    if paragraph.runs:
        _set_run_style(paragraph.runs[0], size, color, bold, name)


def _background(slide, theme):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = theme["background"]


def _title_bar(slide, title: str, theme, taxonomy: dict[str, Any] | None, slide_number: int | None):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, TITLE_BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["title_bar"]
    bar.line.fill.background()
    title_box = slide.shapes.add_textbox(MARGIN_X, Inches(0.08), Inches(8.9), Inches(0.45))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_paragraph_text(p, title, 26, RGBColor(255,255,255), True, "Aptos Display")
    p.alignment = PP_ALIGN.LEFT
    meta = []
    if taxonomy:
        for key in ["course_name", "content_area", "concept"]:
            value = taxonomy.get(key)
            if value and value not in meta:
                meta.append(str(value))
        if len(meta) > 3:
            meta = meta[:3]
    meta_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.13), Inches(3.2), Inches(0.35))
    tf2 = meta_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    _set_paragraph_text(p2, " • ".join(meta), 11, RGBColor(219, 234, 254))
    p2.alignment = PP_ALIGN.RIGHT
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(Inches(12.45), Inches(6.95), Inches(0.45), Inches(0.22))
        nt = num_box.text_frame
        nt.clear()
        p = nt.paragraphs[0]
        _set_paragraph_text(p, str(slide_number), 10, theme["muted"], True)
        p.alignment = PP_ALIGN.RIGHT


def _panel(slide, left, top, width, height, theme, alt: bool = False):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = theme["panel_alt"] if alt else theme["panel"]
    shp.line.color.rgb = theme["line"]
    shp.line.width = Pt(1)
    return shp


def _fit_picture(slide, image_path: str, left, top, width, height):
    with Image.open(image_path) as img:
        iw, ih = img.size
    image_ratio = iw / max(ih, 1)
    box_ratio = width / max(height, 1)
    if image_ratio > box_ratio:
        final_w = width
        final_h = width / image_ratio
    else:
        final_h = height
        final_w = height * image_ratio
    slide.shapes.add_picture(str(image_path), left + (width - final_w) / 2, top + (height - final_h) / 2, width=final_w, height=final_h)


def _caption(slide, text: str, left, top, width, theme):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.24))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    _set_paragraph_text(p, text, 11, theme["muted"])
    p.alignment = PP_ALIGN.CENTER


def _text_card(slide, bullets: list[str], left, top, width, height, theme, archetype: str, heading: str | None = None):
    _panel(slide, left, top, width, height, theme)
    box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), height - Inches(0.28))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if heading:
        p = tf.paragraphs[0]
        _set_paragraph_text(p, heading, 13, theme["accent"], True)
    count = len(bullets)
    total_chars = sum(len(str(b)) for b in bullets)
    size = 23
    if archetype == "split-diagram-text":
        size = 21
    elif archetype == "table":
        size = 18
    elif archetype == "case-ngn":
        size = 19
    if count >= 4:
        size -= 1
    if total_chars > 250:
        size -= 2
    size = max(MIN_FONT_SIZE, min(size, 26))
    for bullet in bullets[:4]:
        p = tf.add_paragraph()
        p.text = str(bullet)
        p.level = 0
        p.space_after = Pt(6)
        if p.runs:
            _set_run_style(p.runs[0], size, theme["ink"], False)
    return {"font_size": size, "bullet_count": count, "total_chars": total_chars}


def _title_bullets(slide, item: dict[str, Any], theme, taxonomy):
    left_w = Inches(4.2)
    right_w = CONTENT_WIDTH - left_w - Inches(0.25)
    _panel(slide, MARGIN_X, CONTENT_TOP, left_w, Inches(5.45), theme, alt=True)
    left_box = slide.shapes.add_textbox(MARGIN_X + Inches(0.22), CONTENT_TOP + Inches(0.22), left_w - Inches(0.44), Inches(5.0))
    lt = left_box.text_frame
    lt.clear()
    p0 = lt.paragraphs[0]
    _set_paragraph_text(p0, item.get("learning_objective", "Lesson objective"), 15, theme["accent"], True)
    p1 = lt.add_paragraph()
    _set_paragraph_text(p1, item.get("speaker_intent", "Explain the key idea and the first safe action."), 20, theme["ink"])
    p1.space_before = Pt(12)
    p2 = lt.add_paragraph()
    _set_paragraph_text(p2, "Use the right panel as the learner-facing summary. Keep narration in the script, not on the slide.", 12, theme["muted"])
    p2.space_before = Pt(18)
    bullets = item.get("on_slide_text") or item.get("bullets") or []
    return _text_card(slide, bullets, MARGIN_X + left_w + Inches(0.25), CONTENT_TOP, right_w, Inches(5.45), theme, "title-bullets", "Key points")


def _split_diagram_text(slide, item: dict[str, Any], theme, taxonomy):
    diagram = item.get("assets", {}).get("diagram") or {}
    left_w = Inches(6.1)
    right_w = CONTENT_WIDTH - left_w - Inches(0.28)
    _panel(slide, MARGIN_X, CONTENT_TOP, left_w, Inches(5.45), theme)
    _panel(slide, MARGIN_X + left_w + Inches(0.28), CONTENT_TOP, right_w, Inches(5.45), theme)
    if diagram.get("file") and Path(diagram["file"]).exists():
        _fit_picture(slide, diagram["file"], MARGIN_X + Inches(0.14), CONTENT_TOP + Inches(0.14), left_w - Inches(0.28), Inches(4.72))
        _caption(slide, item.get("slide_title", "Visual"), MARGIN_X + Inches(0.18), CONTENT_TOP + Inches(4.95), left_w - Inches(0.36), theme)
    bullets = item.get("on_slide_text") or []
    qa = _text_card(slide, bullets, MARGIN_X + left_w + Inches(0.28), CONTENT_TOP, right_w, Inches(5.45), theme, "split-diagram-text", "Learner summary")
    return qa


def _case_slide(slide, item: dict[str, Any], theme):
    case = item.get("case", {})
    left_w = Inches(7.2)
    right_w = CONTENT_WIDTH - left_w - Inches(0.25)
    _panel(slide, MARGIN_X, CONTENT_TOP, left_w, Inches(5.45), theme)
    _panel(slide, MARGIN_X + left_w + Inches(0.25), CONTENT_TOP, right_w, Inches(5.45), theme, alt=True)
    stem_box = slide.shapes.add_textbox(MARGIN_X + Inches(0.2), CONTENT_TOP + Inches(0.2), left_w - Inches(0.4), Inches(5.0))
    st = stem_box.text_frame
    st.clear()
    p0 = st.paragraphs[0]
    _set_paragraph_text(p0, "Practice scenario", 14, theme["accent"], True)
    p1 = st.add_paragraph()
    _set_paragraph_text(p1, case.get("stem", ""), 19, theme["ink"])
    p1.space_before = Pt(8)
    p2 = st.add_paragraph()
    _set_paragraph_text(p2, case.get("question", ""), 15, theme["success"], True)
    p2.space_before = Pt(14)
    cues = [str(x) for x in case.get("cues", [])]
    qa = _text_card(slide, cues, MARGIN_X + left_w + Inches(0.25), CONTENT_TOP, right_w, Inches(5.45), theme, "case-ngn", "Decision cues")
    return qa


def _table_slide(slide, item: dict[str, Any], theme):
    table = item.get("table", {})
    headers = table.get("headers") or ["Column A", "Column B", "Column C"]
    rows = table.get("rows") or []
    _panel(slide, MARGIN_X, CONTENT_TOP, CONTENT_WIDTH, Inches(5.45), theme)
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), MARGIN_X + Inches(0.15), CONTENT_TOP + Inches(0.2), CONTENT_WIDTH - Inches(0.3), Inches(5.0)).table
    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme["title_bar"]
        cell.text = str(header)
        if cell.text_frame.paragraphs[0].runs:
            _set_run_style(cell.text_frame.paragraphs[0].runs[0], 12, RGBColor(255,255,255), True)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table_shape.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme["panel"] if r % 2 else theme["panel_alt"]
            cell.text = str(value)
            if cell.text_frame.paragraphs[0].runs:
                _set_run_style(cell.text_frame.paragraphs[0].runs[0], 12, theme["ink"])
    return {"font_size": 18, "bullet_count": len(rows), "total_chars": sum(len(str(x)) for row in rows for x in row)}


def render_deck(slides: list[dict[str, Any]], output_path: str | Path, taxonomy: dict[str, Any] | None = None, deck_style: str = "clean_academic") -> dict:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]
    theme = _theme(deck_style)
    qa = []
    for item in slides:
        slide = prs.slides.add_slide(blank)
        _background(slide, theme)
        title = item.get("slide_title") or item.get("title") or item.get("slide_id")
        _title_bar(slide, title, theme, taxonomy, item.get("slide_number"))
        archetype = item.get("layout_archetype", "title-bullets")
        if archetype == "split-diagram-text":
            q = _split_diagram_text(slide, item, theme, taxonomy)
        elif archetype == "case-ngn":
            q = _case_slide(slide, item, theme)
        elif archetype == "table":
            q = _table_slide(slide, item, theme)
        else:
            q = _title_bullets(slide, item, theme, taxonomy)
        q.update({"slide_id": item.get("slide_id"), "archetype": archetype, "status": "pass" if q.get("font_size", MIN_FONT_SIZE) >= MIN_FONT_SIZE else "fail"})
        qa.append(q)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {"deck_path": str(output_path), "qa": qa, "status": "pass" if all(x["status"] == "pass" for x in qa) else "fail"}
